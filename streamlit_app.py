import streamlit as st
import openai
from pptx import Presentation
from pptx.util import Inches
from dotenv import load_dotenv
import os
from st_files_connection import FilesConnection
import pandas as pd
import io
import csv
import ast
from datetime import date
import plotly.express as px
from PIL import Image
import plotly.graph_objs as go
import numpy as np

# Load environment variables from a .env file if present
load_dotenv()
path = os.path.dirname(__file__)
# Set your OpenAI API key
openapi_key = st.text_input("OpenAI API Key", type="password")
openai.api_key = openapi_key
today = date.today()

# Title of the Streamlit app
st.title('Slide Content Generator')
def get_franchise_data(topic):
    conn = st.connection('s3', type=FilesConnection)
    df = conn.read("fbc-hackathon-test/growth.csv", input_format="csv", ttl=600) 
    df = df.transpose()
    df.columns = df.iloc[0]  # Use the first row as the header
    df = df.drop(df.index[0])  # Drop the first row since it is now the header
    df = df.to_dict()
    keys_to_keep = topic.split(',')
    keys_to_keep = [key.strip() for key in topic.split(',')]
    filtered_dict = {}
    for data in df:
        if str(data) in keys_to_keep:
            filtered_dict[str(data)] = df[data]

    #st.write(filtered_dict)
    return filtered_dict


def aggr_graph():
    # Example data for multiple franchises
    franchises = ['Franchise A', 'Franchise B', 'Franchise C', 'Franchise D', 'Franchise E']
    sales = [120, 150, 130, 170, 160]

    # Calculate the median value
    median_sales = np.median(sales)

    # Create the bar chart
    fig = go.Figure()

    # Add bars for each franchise
    fig.add_trace(go.Bar(
        x=franchises,
        y=sales,
        name='Sales',
        marker_color='blue'
    ))

    # Add a line representing the median value
    fig.add_trace(go.Scatter(
        x=franchises,
        y=[median_sales] * len(franchises),  # Repeat the median value for each franchise
        mode='lines',
        name='Median Sales',
        line=dict(color='red', dash='dash')
    ))

    # Update layout for better visuals
    fig.update_layout(
        title='Comparative Sales Data Between Franchises and Median Value',
        xaxis_title='Franchise',
        yaxis_title='Sales',
        legend_title='Legend',
        yaxis=dict(showgrid=True),
        xaxis=dict(showgrid=False),
        bargap=0.2  # Gap between bars
    )

    # Show the graph
    fig.show()
    st.plotly_chart(fig)

def add_image(slide, image, left, top, width, height):
    """function to add an image to the PowerPoint slide and specify its position and width"""
    slide.shapes.add_picture(image, left=left, top=top, width=width, height=height)
    
# Function to generate slide content
def generate_aggregate_metrics(content):
    prompt_txt = f"Wait for user input to return a response. Use this data to generate the output as a valid dictionary object:\n\n{str(content)}"
    prompt = f"You are a helpful assistant that generates an executive summary of Franchise's performance metrics. Calculate aggregate metrics for given Franchises and return output a valid dictionary object with key as aggregate_metrics. Do not return anything else."

    # Use ChatCompletion with the new model and API method
    response = openai.chat.completions.create(
        model="gpt-3.5-turbo",  # Specify the model
        messages=[
            {"role": "system", "content": prompt_txt},
            {"role": "user", "content": prompt}
        ],
        temperature=0.7,
    )
    generated_text = response.choices[0].message.content
    #st.write(f"Response: {generated_text}")
    return ast.literal_eval(generated_text)

def generate_key_insights(content):
    prompt_txt = f"Wait for user input to return a response. Use this data to generate the output as a valid dictionary object:\n\n{str(content)}"
    prompt = f"You are a helpful assistant that generates an executive summary of Franchise's performance metrics. Analyze the data and summarize the following trends as brief bullets comparing the performance of the franchises in the enterprise. If the franchise has shown growth, what are the factors contributing to it according to the data given?"
   
    # Use ChatCompletion with the new model and API method
    response = openai.chat.completions.create(
        model="gpt-3.5-turbo",  # Specify the model
        messages=[
            {"role": "system", "content": prompt_txt},
            {"role": "user", "content": prompt}
        ],
        temperature=0.7,
    )
    generated_text = response.choices[0].message.content
    #st.write(f"Response: {generated_text}")
    return generated_text

# function to replace text in pptx first slide with selected filters
def replace_text(replacements, shapes):
    """function to replace text on a PowerPoint slide. Takes dict of {match: replacement, ... } and replaces all matches"""
    for shape in shapes:
        for match, replacement in replacements.items():
            if shape.has_text_frame:
                if (shape.text.find(match)) != -1:
                    text_frame = shape.text_frame
                    for paragraph in text_frame.paragraphs:
                        whole_text = "".join(run.text for run in paragraph.runs)
                        whole_text = whole_text.replace(str(match), str(replacement))
                        for idx, run in enumerate(paragraph.runs):
                            if idx != 0:
                                p = paragraph._p
                                p.remove(run._r)
                        if bool(paragraph.runs):
                            paragraph.runs[0].text = whole_text


# Function to create a PowerPoint presentation
def create_presentation(franchise_data, slide_content, key_insights):
    pptx = path + '//' + 'template.pptx'
    prs = Presentation(pptx)
    bullet_slide_layout = prs.slide_layouts.get_by_name('Purple_Speaker1')
    insights_slide_layout = prs.slide_layouts.get_by_name('Purple_Circle_Corners')

    details_dict = {
    'Franchisee': 'Franchisee',
    'NetworkPerformancePartner': 'FBC',
    'Region': 'DO',
    'WeightedScore': 'Your Total Score',
    'Rank': 'Network Ranking',
    'CurrentYearTotalBillableHours': 'Current Year Billable Hours',
    'LastYearTotalBillableHours' : 'Previous Year Billable Hours', 
    'CurrentYearTotalRevenue': 'Current Year Total Revenue',
    'LastYearYearTotalRevenue' : 'Previous Year Total Revenue',     
    'HoursGrowth': 'Your Hours Growth %',
    'RevenueGrowth': 'Your Revenue Growth %',
    'HoursGrowth': 'Your Hours Growth %', 
    'RevenueGrowth': 'Your Revenue Growth %'
    }

    median_data = {"RPNLeadsGrowth": 14.22, "HoursGrowth": 3.97, "RevenueGrowth": 10.23}


    owner = []
    owner_full_name = []
    aggregate_metrics = {}
    if 'aggregate_metrics' in slide_content:
        aggregate_metrics = slide_content['aggregate_metrics']
    if 'AggregateMetrics' in slide_content:
        aggregate_metrics = slide_content['AggregateMetrics']
    #st.write(aggregate_metrics)

    franchise_numbers = []
    
    for franchise in franchise_data:
        franchise_numbers.append(str(franchise))
        slide = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]
        #perf_shape = slide.placeholders[2]
        tf = body_shape.text_frame
        ind_fran = franchise_data[str(franchise)]
        owner_full_name.append(ind_fran['FirstName'] + ' ' + ind_fran['LastName']) 
        owner.append(ind_fran['LastName'])
        title_shape.text = f"Franchise {franchise} - {ind_fran['FirstName']} {ind_fran['LastName']}"
        p = tf.add_paragraph()
        weighted_score =  ind_fran['WeightedScore']
        p.text+= f"FBC: {ind_fran['NetworkPerformancePartner']}\n\n"
        p.text+= f"DO: {ind_fran['Region']}\n\n"
        p.text+= f"Your Score: {weighted_score}\n\n"
        p.text+= f"Rank: {ind_fran['Rank']}\n\n"
        performance_standing = ''
        if weighted_score >=0 and weighted_score < 1.99:
            performance_standing = "Significantly Below Target"
        elif weighted_score >=1.99 and weighted_score < 2.99:
            performance_standing = "Below Target"
        elif weighted_score >=2.99 and weighted_score < 3.99:
            performance_standing = "On Target"
        elif weighted_score >=3.99 and weighted_score < 4.99:
            performance_standing = "Above Target"
        elif weighted_score >=4.99:
            performance_standing = "Significantly Above Target"
        p.text+= f"Performance Standing: {performance_standing}\n\n"
        for placeholder in shapes.placeholders:
            #st.write(placeholder.name)
            if placeholder.name == 'Picture Placeholder 1':
                df = pd.DataFrame(
                    [[str(franchise), float(ind_fran['RevenueGrowth']), float(ind_fran['HoursGrowth']), float(ind_fran['RPNLeadsGrowth'])], 
                    ["N/W Median", float(median_data['RevenueGrowth']), float(median_data['HoursGrowth']), float(median_data['RPNLeadsGrowth'])]],
                columns=["Franchise", "Revenue", "Billable Hours", "RPN Leads"]
                )

                width = Inches(8)
                left = Inches(2.5)
                top = Inches(1)
                fig = px.bar(df, x="Franchise", y=["Revenue", "Billable Hours", "RPN Leads"], barmode='group', height=400)
                fig.update_layout(paper_bgcolor='rgba(0,0,0,0)', plot_bgcolor='rgba(0,0,0,0)')
                # st.dataframe(df) # if need to display dataframe
                #st.plotly_chart(fig)
                metrics_im = 'metrics.png'

                fig.write_image(metrics_im)
                with Image.open(metrics_im) as img:
                    image_width, image_height = img.size
                placeholder_width = placeholder.width
                placeholder_height = placeholder.height
                left = placeholder.left
                top = placeholder.top
                
                # Calculate aspect ratios
                image_ratio = image_width / image_height
                placeholder_ratio = placeholder_width / placeholder_height

                # Determine the scaling factor
                if image_ratio > placeholder_ratio:
                    # Image is wider than placeholder
                    scale_factor = placeholder_width / image_width
                else:
                    # Image is taller than placeholder
                    scale_factor = placeholder_height / image_height

                # Calculate new dimensions for the image
                new_width = int(image_width * scale_factor)
                new_height = int(image_height * scale_factor)    
                
                add_image(slide, image=metrics_im, left=left, width=new_width, height=new_height, top=top)
                os.remove('metrics.png')
        #for k in ind_fran:
        #    if k in details_dict:
        #        p = tf.add_paragraph()
        #        p.text+= f"  {details_dict[k]}: {ind_fran[k]}\n\n"


    franchise_numbers_string = ", ".join(franchise_numbers)
    owner = list(set(owner))
    #st.write(owner)
    # Convert the array to a comma-separated string
    owner_string = ", ".join(owner)
    #st.write(f"comma separated unique list: {owner_string}")

    owner_full_name = list(set(owner_full_name))
    owner_full_name_string = ", ".join(owner_full_name)
    #st.write(f"comma separated unique list: {owner_full_name_string}")
    
    first_slide = prs.slides[0]
    shapes_1 = []

    third_slide = prs.slides[2]
    shapes_2 = []

    #Aggregate Metrics
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    title_shape.text = f"Enterprise Business Overview"
    #body_shape = shapes.placeholders[1]
    #tf = body_shape.text_frame
    #for k in aggregate_metrics:
    #    p = tf.add_paragraph()
    #    p.text+= f"  {k}: {aggregate_metrics[k]}\n\n"

    aggr_graph()

    #Key Insights
    slide = prs.slides.add_slide(insights_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    title_shape.text = f"Key Insights"
    body_shape = shapes.placeholders[1]
    tf = body_shape.text_frame
    p = tf.add_paragraph()
    p.text+= key_insights
        
    # create lists with shape objects
    for shape in first_slide.shapes:
        shapes_1.append(shape)

    # create lists with shape objects
    for shape in third_slide.shapes:
        shapes_2.append(shape)

    # initiate a dictionary of placeholders and values to replace
    replaces_1 = {
        '{owner}': owner_string,
        '{date}': today }
    replace_text(replaces_1, shapes_1)

    replaces_2 = {
        '{owner_name}': owner_full_name_string,
        '{franchise_numbers}': franchise_numbers_string }
    replace_text(replaces_2, shapes_2)

    # Save the presentation
    file_path = "generated_presentation.pptx"
    prs.save(file_path)
    return file_path

# Streamlit input fields
topic = st.text_input("Enter the Franchise number:")
content = st.text_area("Enter the themes for the slides:")

# Generate button
if st.button("Generate Slide Content"):
    if topic:
        franchise_data = get_franchise_data(topic)
        st.write(franchise_data)
        generated_content = generate_aggregate_metrics(franchise_data)
        
        key_insights = generate_key_insights(franchise_data)
        st.subheader("Generated Slide Content:")
        
        # Create and offer download of the PowerPoint presentation
        file_path = create_presentation(franchise_data, generated_content, key_insights)
        
        with open(file_path, "rb") as file:
            btn = st.download_button(
                label="Download PowerPoint Presentation",
                data=file,
                file_name=file_path,
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
            )
    else:
        st.error("Please enter both the topic and content.") 
